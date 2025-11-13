from pptx import Presentation

def load_presentation(file_path):
    """
    Extract text from all slides and return structured JSON.
    """
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if shape == slide.shapes[0]:
                    title = text
                else:
                    slide_content.append(text)
        slides_text.append({
            "slide_number": i + 1,
            "title": title,
            "content": slide_content
        })
    return slides_text